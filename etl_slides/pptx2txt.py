import logging
import glob
from tqdm import tqdm
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)


def pptx2txt():
    logging.info("Starting pptx2txt conversion...")
    files = glob.glob("slides_adapt\\*.pptx")
    logging.info(f"Files found: {files}")

    for file in files:
        logging.info(f"Processing file: {file}")
        new_file = "slides_txt\\" + file.split("\\")[1] + ".txt"

        f = open(new_file, "w", encoding="utf-8")

        prs = Presentation(file)
        for i, slide in tqdm(
            enumerate(prs.slides), desc="Processing slides", unit="slide"
        ):

            f.write(f"\nSLIDE_{i}\n")

            for shape in slide.shapes:
                f.write(f"SHAPE_{shape.shape_type}")

                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    f.write(f"_AUTO_SHAPE_{shape.auto_shape_type}\n")
                    f.write(f"HEIGHT_{shape.height}\n")
                    f.write(f"TOP_{shape.top}\n")

                if hasattr(shape, "text"):
                    f.write(f"\nSTART_TEXT\n{shape.text}\nEND_TEXT\n")

                f.write("\n")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.lower() == "índice":
                        f.write("\n__END__\n")

        f.close()


def main():
    pptx2txt()


if __name__ == "__main__":
    main()
