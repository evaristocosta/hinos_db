import glob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def main():
    files = glob.glob("slides_adapt\\*.pptx")

    for file in files:
        new_file = "slides_txt\\" + file.split("\\")[1] + ".txt"

        f = open(new_file, "w")

        prs = Presentation(file)
        for i, slide in enumerate(prs.slides):

            f.write(f"\nSLIDE_{i}\n")

            for shape in slide.shapes:
                f.write(f"SHAPE_{shape.shape_type}")

                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    f.write(f"_AUTO_SHAPE_{shape.auto_shape_type}\n")
                    f.write(f"HEIGHT_{shape.height}\n")
                    f.write(f"TOP_{shape.top}\n")

                if hasattr(shape, "text"):
                    f.write(f"\nSTART_TEXT\n{shape.text}\nEND_TEXT\n")

                f.write("\n")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.lower() == "índice":
                        f.write("\n__END__\n")

        f.close()


if __name__ == "__main__":
    main()
